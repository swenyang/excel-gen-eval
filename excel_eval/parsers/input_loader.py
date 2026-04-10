"""Load grounding data from various file formats."""

from __future__ import annotations

import json
from pathlib import Path

import pandas as pd


def load_input_file(path: str | Path) -> str:
    """Load a single input file and return its content as text.

    Supports: CSV, JSON, JSONL, Excel (.xlsx/.xls), plain text.
    For binary/structured formats, converts to a readable text representation.
    """
    path = Path(path)
    if not path.exists():
        raise FileNotFoundError(f"Input file not found: {path}")

    suffix = path.suffix.lower()

    if suffix == ".csv":
        return _load_csv(path)
    elif suffix == ".json":
        return _load_json(path)
    elif suffix == ".jsonl":
        return _load_jsonl(path)
    elif suffix in (".xlsx", ".xls"):
        return _load_excel(path)
    elif suffix == ".pdf":
        return _load_pdf(path)
    elif suffix in (".docx", ".doc"):
        return _load_docx(path)
    elif suffix == ".pptx":
        return _load_pptx(path)
    elif suffix in (".jpg", ".jpeg", ".png", ".gif", ".webp"):
        return f"[Image file: {path.name} — visual content not extractable as text]"
    else:
        return _load_text(path)


def load_all_inputs(file_configs: list[dict], base_dir: str | Path) -> str:
    """Load all input files and combine into a single text block.

    Each file_config should have 'path' and optionally 'description'.
    Returns a combined text with file headers.
    """
    base_dir = Path(base_dir)
    parts: list[str] = []

    for fc in file_configs:
        file_path = base_dir / fc["path"]
        desc = fc.get("description", "")
        header = f"--- Input File: {fc['path']}"
        if desc:
            header += f" ({desc})"
        header += " ---"

        try:
            content = load_input_file(file_path)
            parts.append(f"{header}\n{content}")
        except Exception as e:
            parts.append(f"{header}\n[Error loading file: {e}]")

    return "\n\n".join(parts)


def _load_csv(path: Path) -> str:
    """Load CSV file, return as CSV text."""
    return path.read_text(encoding="utf-8")


def _load_json(path: Path) -> str:
    """Load JSON file, return as pretty-printed JSON."""
    with open(path, "r", encoding="utf-8") as f:
        data = json.load(f)
    return json.dumps(data, indent=2, ensure_ascii=False)


def _load_jsonl(path: Path) -> str:
    """Load JSONL file, return each line pretty-printed."""
    lines = path.read_text(encoding="utf-8").strip().split("\n")
    formatted = []
    for i, line in enumerate(lines):
        try:
            obj = json.loads(line)
            formatted.append(f"Record {i + 1}: {json.dumps(obj, ensure_ascii=False)}")
        except json.JSONDecodeError:
            formatted.append(f"Record {i + 1}: {line}")
    return "\n".join(formatted)


def _load_excel(path: Path) -> str:
    """Load Excel file as input data, return CSV-like text per sheet.

    Uses the same display-value formatting as excel_parser to ensure
    source and generated data use identical value representations.
    """
    from .excel_parser import parse_excel

    try:
        parsed = parse_excel(str(path))
        parts: list[str] = []
        for sheet in parsed.sheets:
            parts.append(f"[Sheet: {sheet.name}]\n{sheet.csv_text}")
        return "\n\n".join(parts)
    except Exception:
        # Fallback to pandas if excel_parser fails
        parts = []
        xls = pd.ExcelFile(path)
        for sheet_name in xls.sheet_names:
            df = pd.read_excel(xls, sheet_name=sheet_name)
            csv_text = df.to_csv(index=False)
            parts.append(f"[Sheet: {sheet_name}]\n{csv_text}")
        return "\n\n".join(parts)


def _load_text(path: Path) -> str:
    """Load plain text file."""
    try:
        return path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        return path.read_text(encoding="latin-1")


def _load_pdf(path: Path) -> str:
    """Load PDF file, extracting text and tables via pdfplumber."""
    try:
        import pdfplumber
    except ImportError:
        return f"[PDF file: {path.name} — install pdfplumber for better extraction]"

    parts: list[str] = []
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages):
            page_parts: list[str] = []

            text = page.extract_text()
            if text:
                page_parts.append(text)

            tables = page.extract_tables()
            for t_idx, table in enumerate(tables):
                if table:
                    rows = [" | ".join(str(cell or "") for cell in row) for row in table]
                    page_parts.append(f"[Table {t_idx + 1}]\n" + "\n".join(rows))

            if page_parts:
                parts.append(f"[Page {i + 1}]\n" + "\n".join(page_parts))

    return "\n\n".join(parts) if parts else f"[PDF file: {path.name} — no extractable text]"


def _load_docx(path: Path) -> str:
    """Load Word document, extracting paragraphs and tables."""
    try:
        from docx import Document
    except ImportError:
        return f"[Word file: {path.name} — install python-docx for extraction]"

    doc = Document(str(path))
    parts: list[str] = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            parts.append(text)

    for i, table in enumerate(doc.tables):
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        if rows:
            parts.append(f"[Table {i + 1}]\n" + "\n".join(rows))

    return "\n\n".join(parts) if parts else f"[Word file: {path.name} — no extractable text]"


def _load_pptx(path: Path) -> str:
    """Load PowerPoint file, extracting slide text."""
    try:
        from pptx import Presentation
    except ImportError:
        return f"[PowerPoint file: {path.name} — install python-pptx for extraction]"

    prs = Presentation(str(path))
    parts: list[str] = []

    for i, slide in enumerate(prs.slides):
        slide_text: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_text.append(text)
            if shape.has_table:
                rows = []
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append(" | ".join(cells))
                if rows:
                    slide_text.append("[Table]\n" + "\n".join(rows))
        if slide_text:
            parts.append(f"[Slide {i + 1}]\n" + "\n".join(slide_text))

    return "\n\n".join(parts) if parts else f"[PowerPoint file: {path.name} — no extractable text]"
